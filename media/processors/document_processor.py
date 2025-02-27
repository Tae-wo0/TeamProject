import os
import io
import json
import logging
import requests
from typing import Dict, List
from urllib.parse import urlparse
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from openpyxl import load_workbook
import openai
from config import OPENAI_API_KEY
import chardet

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.StreamHandler()],
)

openai.api_key = OPENAI_API_KEY


class DocumentProcessor:
    def __init__(self, base_manager):
        self.base_manager = base_manager

    def extract_pdf_content(self, file_obj) -> Dict[str, str]:
        try:
            logging.info(" PDF 파일 처리 시작")
            reader = PdfReader(file_obj)
            title = reader.metadata.title if reader.metadata and reader.metadata.title else "PDF Document"
            content = "\n".join([page.extract_text() or "" for page in reader.pages]).strip()
            return {"title": title, "content": content} if content else None
        except Exception as e:
            logging.error(f" PDF 처리 실패: {str(e)}")
            return None

    def extract_docx_content(self, file_obj) -> Dict[str, str]:
        try:
            doc = docx.Document(file_obj)
            content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            return {"title": "Word Document", "content": content}
        except Exception as e:
            logging.error(f" DOCX 처리 실패: {str(e)}")
            return None

    def extract_pptx_content(self, file_obj) -> Dict[str, str]:
        try:
            prs = Presentation(file_obj)
            content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            return {"title": "PowerPoint Document", "content": content}
        except Exception as e:
            logging.error(f" PPTX 처리 실패: {str(e)}")
            return None

    def extract_xlsx_content(self, file_obj) -> Dict[str, str]:
        try:
            wb = load_workbook(file_obj)
            content = "\n".join(
                f"[{sheet}]\n" + "\n".join(
                    " | ".join(str(cell) for cell in row if cell) for row in wb[sheet].iter_rows(values_only=True)
                ) for sheet in wb.sheetnames
            ).strip()
            return {"title": "Excel Document", "content": content}
        except Exception as e:
            logging.error(f" Excel 처리 실패: {str(e)}")
            return None

    def create_chunks(self, text: str, chunk_size: int = 300, overlap_size: int = 75) -> List[str]:
        """텍스트를 300자 청크 + 75자 오버랩 방식으로 분할"""
        chunks = []
        words = text.split()
        total_length = len(words)

        if total_length < chunk_size:
            return [text]  # 작은 텍스트는 그대로 반환

        for i in range(0, total_length, chunk_size - overlap_size):
            chunk = " ".join(words[i:i + chunk_size])
            if len(chunk) > chunk_size * 0.3:  # 최소 크기 확인
                chunks.append(chunk)

        logging.info(f" 생성된 청크 수: {len(chunks)}")
        return chunks

    def process_document_url(self, file_url: str):
        """문서 URL을 처리하여 300자 청크 + 75자 오버랩 방식으로 변환"""
        try:
            logging.info(f"문서 다운로드 및 처리 시작: {file_url}")
            parsed_url = urlparse(file_url)
            file_ext = os.path.splitext(parsed_url.path)[1].lower().replace(".", "")

            response = requests.get(file_url)
            if response.status_code != 200:
                raise Exception(f" 문서 다운로드 실패: {response.status_code}")

            file_content = io.BytesIO(response.content)
            extracted_data = None

            if file_ext == "pdf":
                extracted_data = self.extract_pdf_content(file_content)
            elif file_ext == "docx":
                extracted_data = self.extract_docx_content(file_content)
            elif file_ext == "pptx":
                extracted_data = self.extract_pptx_content(file_content)
            elif file_ext == "xlsx":
                extracted_data = self.extract_xlsx_content(file_content)
            elif file_ext == "txt":
                raw_data = response.content
                detected_encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
                content_text = raw_data.decode(detected_encoding, errors='ignore').strip()
                extracted_data = {"title": "Text File", "content": content_text}
            else:
                raise Exception(f" 지원하지 않는 문서 형식: {file_ext}")

            if not extracted_data or not extracted_data.get('content', '').strip():
                raise Exception(" 문서에서 텍스트를 추출할 수 없음!")

            # 300자 청크 + 75자 오버랩 적용
            content = extracted_data['content']
            extracted_data['chunks'] = self.create_chunks(content, chunk_size=300, overlap_size=75)

            logging.info(f" 문서 제목: {extracted_data['title']}")
            logging.info(f" 첫 번째 청크 (300자 기준): {extracted_data['chunks'][0][:300]}...")
            return extracted_data

        except Exception as e:
            logging.error(f" 문서 처리 중 오류 발생: {str(e)}")
            return {"success": False, "message": str(e)}
